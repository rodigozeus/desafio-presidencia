"""
Gera o arquivo PowerPoint da apresentação do desafio PR-CGSE.
Uso: python scripts/gerar_apresentacao.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Paleta ──────────────────────────────────────────────────────────────────
VERDE_GOV    = RGBColor(0x15, 0x63, 0x4E)   # verde escuro institucional
VERDE_CLARO  = RGBColor(0x2D, 0xA7, 0x71)   # verde destaque
CINZA_ESCURO = RGBColor(0x2B, 0x2D, 0x42)   # texto principal
CINZA_MEDIO  = RGBColor(0x55, 0x59, 0x6E)   # texto secundário
BRANCO       = RGBColor(0xFF, 0xFF, 0xFF)
AMARELO      = RGBColor(0xFF, 0xC0, 0x00)
AZUL_CLARO   = RGBColor(0xE8, 0xF4, 0xF1)   # fundo de caixas

# ── Helpers ──────────────────────────────────────────────────────────────────
def add_text_box(slide, text, left, top, width, height,
                 font_size=18, bold=False, color=None, align=PP_ALIGN.LEFT,
                 bg_color=None, italic=False, wrap=True):
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color or CINZA_ESCURO
    if bg_color:
        fill = txBox.fill
        fill.solid()
        fill.fore_color.rgb = bg_color
    return txBox

def add_rect(slide, left, top, width, height, fill_color, line_color=None):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    return shape

def slide_header(slide, title, subtitle=None, accent=VERDE_GOV):
    """Faixa superior colorida + título."""
    add_rect(slide, 0, 0, 13.33, 1.4, accent)
    add_text_box(slide, title, 0.4, 0.1, 12.5, 0.9,
                 font_size=28, bold=True, color=BRANCO, align=PP_ALIGN.LEFT)
    if subtitle:
        add_text_box(slide, subtitle, 0.4, 0.95, 12.5, 0.4,
                     font_size=14, color=AMARELO, align=PP_ALIGN.LEFT)

def add_bullet_box(slide, items, left, top, width, height,
                   title=None, font_size=15, icon="▸"):
    if title:
        add_text_box(slide, title, left, top, width, 0.45,
                     font_size=16, bold=True, color=VERDE_GOV)
        top += 0.45
        height -= 0.45
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = f"{icon}  {item}"
        run.font.size = Pt(font_size)
        run.font.color.rgb = CINZA_ESCURO
    return txBox

def add_code_box(slide, code, left, top, width, height, font_size=11):
    rect = add_rect(slide, left, top, width, height,
                    RGBColor(0x1E, 0x1E, 0x2E), RGBColor(0x44, 0x47, 0x5A))
    txBox = slide.shapes.add_textbox(
        Inches(left + 0.15), Inches(top + 0.15),
        Inches(width - 0.3), Inches(height - 0.3)
    )
    tf = txBox.text_frame
    tf.word_wrap = False
    first = True
    for line in code.split("\n"):
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        run = p.add_run()
        run.text = line
        run.font.size = Pt(font_size)
        run.font.color.rgb = RGBColor(0xCD, 0xD6, 0xF4)
        run.font.name = "Courier New"
    return txBox

def badge(slide, text, left, top, width=1.6, height=0.38,
          bg=VERDE_CLARO, fg=BRANCO, font_size=12):
    add_rect(slide, left, top, width, height, bg)
    add_text_box(slide, text, left, top, width, height,
                 font_size=font_size, bold=True, color=fg,
                 align=PP_ALIGN.CENTER)


# ── Criação da apresentação ──────────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

blank_layout = prs.slide_layouts[6]   # completamente em branco


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — CAPA
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)

add_rect(sl, 0, 0, 13.33, 7.5, VERDE_GOV)               # fundo
add_rect(sl, 0, 5.6, 13.33, 1.9, RGBColor(0x0D, 0x3D, 0x2F))  # rodapé

add_text_box(sl, "Presidência da República", 1.0, 0.6, 11.33, 0.6,
             font_size=16, color=AMARELO, align=PP_ALIGN.CENTER)
add_text_box(sl, "Plataforma de Gestão de Pedidos", 1.0, 1.3, 11.33, 1.2,
             font_size=38, bold=True, color=BRANCO, align=PP_ALIGN.CENTER)
add_text_box(sl, "Produto Mínimo Viável — Arquitetura de Microsserviços + Microfrontends",
             1.0, 2.55, 11.33, 0.7,
             font_size=18, color=RGBColor(0xB5, 0xE8, 0xD5), align=PP_ALIGN.CENTER)

add_text_box(sl, "Edital Nº 173/2026  ·  PR-CGSE", 1.0, 5.75, 11.33, 0.5,
             font_size=14, color=RGBColor(0xB5, 0xE8, 0xD5), align=PP_ALIGN.CENTER)

# badges tecnologias
bx = 2.2; by = 3.6; bw = 1.55; bh = 0.42; gap = 1.75
for txt, col in [
    ("FastAPI",    RGBColor(0x00, 0x92, 0x68)),
    ("PostgreSQL", RGBColor(0x33, 0x6B, 0x91)),
    ("Redis",      RGBColor(0xC0, 0x30, 0x2A)),
    ("React MFE",  RGBColor(0x61, 0x4F, 0xC7)),
    ("Claude IA",  RGBColor(0xCC, 0x7A, 0x00)),
]:
    badge(sl, txt, bx, by, bw, bh, bg=col)
    bx += gap


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — AGENDA
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
add_rect(sl, 0, 0, 13.33, 7.5, RGBColor(0xF4, 0xF7, 0xF6))

slide_header(sl, "Agenda da Apresentação", "20 – 30 minutos")

items = [
    ("1. Visão Geral da Arquitetura",     "Serviços, comunicação, MFEs e banco de dados",          2.0),
    ("2. Decisões Técnicas",               "FastAPI vs Django, Redis, JWT, Module Federation",      3.1),
    ("3. Demonstração ao Vivo",            "docker-compose up → criar pedido → filtrar → status",   4.2),
    ("4. Tour pelo Código",                "Separação de serviços, MFE, IA integrada",              5.3),
    ("5. Decisões Não Tomadas",            "O que ficou de fora e o que priorizaria a seguir",      6.1),
]

for titulo, desc, top in items:
    add_rect(sl, 0.6, top - 0.05, 12.1, 0.78, BRANCO, VERDE_CLARO)
    add_text_box(sl, titulo, 0.8, top,     6.0, 0.45, font_size=16, bold=True,  color=VERDE_GOV)
    add_text_box(sl, desc,   6.9, top,     5.8, 0.45, font_size=14, color=CINZA_MEDIO)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — CONTEXTO DO DESAFIO
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
add_rect(sl, 0, 0, 13.33, 7.5, RGBColor(0xF4, 0xF7, 0xF6))
slide_header(sl, "Contexto do Desafio", "Problema que estamos resolvendo")

# Problema
add_rect(sl, 0.5, 1.65, 5.7, 4.8, BRANCO, RGBColor(0xCC, 0x30, 0x2A))
add_text_box(sl, "😟  Situação Atual", 0.7, 1.75, 5.2, 0.5, font_size=16, bold=True, color=RGBColor(0xCC, 0x30, 0x2A))
add_bullet_box(sl, [
    "Controle de pedidos em planilhas",
    "Processo lento e propenso a erros",
    "Sem visibilidade em tempo real",
    "Sem rastreamento de status",
    "Sem histórico estruturado",
], 0.7, 2.3, 5.2, 4.0, font_size=15)

# Solução
add_rect(sl, 6.9, 1.65, 5.9, 4.8, BRANCO, VERDE_CLARO)
add_text_box(sl, "✅  PMV Entregue", 7.1, 1.75, 5.4, 0.5, font_size=16, bold=True, color=VERDE_GOV)
add_bullet_box(sl, [
    "Plataforma web com autenticação",
    "Criação e gestão de pedidos",
    "Filtro por status em tempo real",
    "Sugestão de prioridade com IA",
    "Arquitetura escalável e independente",
], 7.1, 2.3, 5.4, 4.0, font_size=15)

add_text_box(sl, "→", 6.2, 3.6, 0.7, 0.7, font_size=36, bold=True, color=VERDE_CLARO, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — ARQUITETURA GERAL
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
add_rect(sl, 0, 0, 13.33, 7.5, RGBColor(0xF4, 0xF7, 0xF6))
slide_header(sl, "Visão Geral da Arquitetura", "5 serviços independentes + 2 bancos + cache")

# Camada Frontend
add_rect(sl, 0.4, 1.6, 12.5, 1.7, RGBColor(0xEB, 0xF5, 0xFF), RGBColor(0x61, 0x4F, 0xC7))
add_text_box(sl, "FRONTEND  (porta 3000)", 0.6, 1.65, 4.0, 0.4, font_size=13, bold=True, color=RGBColor(0x61, 0x4F, 0xC7))

badge(sl, "MFE Shell / Host", 0.7, 2.05, 2.8, 0.6, bg=RGBColor(0x61, 0x4F, 0xC7))
add_text_box(sl, "Login · Navegação · AuthContext", 0.7, 2.7, 2.8, 0.35, font_size=10, color=CINZA_MEDIO, align=PP_ALIGN.CENTER)

add_text_box(sl, "◀── Module Federation ──▶", 3.7, 2.2, 3.2, 0.5, font_size=12, color=CINZA_MEDIO, align=PP_ALIGN.CENTER)

badge(sl, "MFE Orders  :3001", 7.0, 2.05, 2.8, 0.6, bg=RGBColor(0x41, 0x3A, 0x9A))
add_text_box(sl, "Lista · Criar · Detalhe", 7.0, 2.7, 2.8, 0.35, font_size=10, color=CINZA_MEDIO, align=PP_ALIGN.CENTER)

# Setas para backend
add_text_box(sl, "HTTP + JWT Bearer", 1.5, 3.45, 2.5, 0.4, font_size=10, color=CINZA_MEDIO, align=PP_ALIGN.CENTER)
add_text_box(sl, "↓", 2.4, 3.3, 0.8, 0.5, font_size=22, bold=True, color=CINZA_MEDIO, align=PP_ALIGN.CENTER)
add_text_box(sl, "HTTP + JWT Bearer", 7.7, 3.45, 2.5, 0.4, font_size=10, color=CINZA_MEDIO, align=PP_ALIGN.CENTER)
add_text_box(sl, "↓", 8.6, 3.3, 0.8, 0.5, font_size=22, bold=True, color=CINZA_MEDIO, align=PP_ALIGN.CENTER)

# Camada backend
badge(sl, "Users Service  :8001", 0.6, 3.85, 3.2, 0.62, bg=VERDE_GOV)
add_text_box(sl, "FastAPI · PostgreSQL\nJWT Issuer · /auth /users", 0.6, 4.52, 3.2, 0.6, font_size=10, color=CINZA_MEDIO, align=PP_ALIGN.CENTER)

badge(sl, "Orders Service  :8002", 7.0, 3.85, 3.2, 0.62, bg=VERDE_GOV)
add_text_box(sl, "FastAPI · PostgreSQL · Redis\nJWT Validator · /orders + IA", 7.0, 4.52, 3.2, 0.6, font_size=10, color=CINZA_MEDIO, align=PP_ALIGN.CENTER)

# Bancos
badge(sl, "users_db\nPostgreSQL", 0.6, 5.5, 2.2, 0.75, bg=RGBColor(0x33, 0x6B, 0x91))
badge(sl, "orders_db\nPostgreSQL", 6.8, 5.5, 2.2, 0.75, bg=RGBColor(0x33, 0x6B, 0x91))
badge(sl, "Redis Cache\n60s TTL", 9.3, 5.5, 2.0, 0.75, bg=RGBColor(0xC0, 0x30, 0x2A))

# Claude API (externo)
badge(sl, "Claude API (IA)\nAnthropic", 9.8, 3.85, 2.8, 0.62, bg=RGBColor(0xCC, 0x7A, 0x00))
add_text_box(sl, "↖ chama ao criar pedido", 9.5, 4.52, 3.0, 0.35, font_size=9, color=CINZA_MEDIO)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — FLUXO DE AUTENTICAÇÃO JWT
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
add_rect(sl, 0, 0, 13.33, 7.5, RGBColor(0xF4, 0xF7, 0xF6))
slide_header(sl, "Autenticação JWT Compartilhada", "Shared secret — validação local em cada serviço")

steps = [
    ("1", "Usuário faz login",           "Frontend envia e-mail + senha ao\nUsers Service (POST /auth/login)"),
    ("2", "JWT é emitido",               "Users Service valida credenciais,\ngera token HS256 (24h)"),
    ("3", "Token no localStorage",       "Frontend armazena e injeta como\nAuthorization: Bearer em cada request"),
    ("4", "Orders Service valida",       "Verifica assinatura com JWT_SECRET\ncompartilhado — sem chamada inter-serviço"),
]

for i, (num, titulo, desc) in enumerate(steps):
    col = 0.5 + i * 3.2
    add_rect(sl, col, 1.7, 2.9, 3.8, BRANCO, VERDE_CLARO)
    add_text_box(sl, num, col + 0.1, 1.75, 0.7, 0.7,
                 font_size=28, bold=True, color=VERDE_CLARO)
    add_text_box(sl, titulo, col + 0.15, 2.55, 2.6, 0.55,
                 font_size=14, bold=True, color=CINZA_ESCURO)
    add_text_box(sl, desc, col + 0.15, 3.15, 2.6, 1.4,
                 font_size=12, color=CINZA_MEDIO)
    if i < 3:
        add_text_box(sl, "→", col + 2.95, 2.8, 0.35, 0.5,
                     font_size=24, bold=True, color=VERDE_CLARO)

add_rect(sl, 0.5, 5.8, 12.3, 1.35, AZUL_CLARO, VERDE_CLARO)
add_text_box(sl, "⚠  Decisão de design documentada:",
             0.7, 5.85, 5.0, 0.4, font_size=13, bold=True, color=VERDE_GOV)
add_text_box(sl,
    "Alternativa não tomada: JWKS endpoint para verificação criptográfica "
    "— mais robusto (permite revogação imediata), mas adiciona latência em cada request. "
    "Para o PMV, shared secret é suficiente.",
    0.7, 6.3, 12.0, 0.7, font_size=12, color=CINZA_ESCURO)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — DECISÕES TÉCNICAS
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
add_rect(sl, 0, 0, 13.33, 7.5, RGBColor(0xF4, 0xF7, 0xF6))
slide_header(sl, "Decisões Técnicas", "Por que cada tecnologia foi escolhida")

decisoes = [
    ("FastAPI\n(não Django)",
     ["Async-native, sem overhead do Django ORM",
      "Swagger/OpenAPI automático em /docs",
      "Pydantic v2 para validação com type hints",
      "Menos boilerplate para um PMV"]),
    ("Webpack Module\nFederation (MFE)",
     ["Bundle separado carregado em runtime",
      "Deploy independente sem rebuild do Shell",
      "React + Router compartilhados como singletons",
      "Alternativa rejeitada: iframes (UX ruim)"]),
    ("Redis como cache",
     ["TTL 60s na listagem de pedidos",
      "Invalidado ao criar/atualizar pedido",
      "Falha graciosamente: API continua sem cache",
      "Mais simples que MongoDB para este caso"]),
    ("Claude Haiku (IA)",
     ["Sugere prioridade + gera resumo do pedido",
      "Haiku: baixa latência e custo reduzido",
      "Opcional: sem API key usa fallback padrão",
      "Prompt estruturado com regras de negócio"]),
]

for i, (titulo, bullets) in enumerate(decisoes):
    col = 0.4 + (i % 2) * 6.5
    row = 1.7 + (i // 2) * 2.7
    add_rect(sl, col, row, 6.1, 2.4, BRANCO, VERDE_GOV)
    add_text_box(sl, titulo, col + 0.2, row + 0.1, 5.7, 0.7,
                 font_size=15, bold=True, color=VERDE_GOV)
    add_bullet_box(sl, bullets, col + 0.2, row + 0.82, 5.7, 1.5,
                   font_size=12, icon="·")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — MICROFRONTENDS (MFE)
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
add_rect(sl, 0, 0, 13.33, 7.5, RGBColor(0xF4, 0xF7, 0xF6))
slide_header(sl, "Estratégia de Microfrontends", "Webpack 5 Module Federation")

# diagrama textual
add_rect(sl, 0.4, 1.65, 12.4, 2.5, RGBColor(0x1E, 0x1E, 0x2E))
code_arch = """\
  Shell (Host) :3000                       MFE Orders (Remote) :3001
  ┌─────────────────────────────┐          ┌───────────────────────────┐
  │  webpack.config.js          │          │  webpack.config.js        │
  │  remotes: {                 │  carrega │  exposes: {               │
  │    mfe_orders: ":3001/      │ ────────►│    './OrdersApp': ...     │
  │      remoteEntry.js"        │  runtime │  }                        │
  │  }                          │          │                           │
  └─────────────────────────────┘          └───────────────────────────┘"""
add_code_box(sl, code_arch, 0.5, 1.7, 12.2, 2.35, font_size=11)

add_bullet_box(sl, [
    "Shell é o HOST: define as rotas, provê AuthContext e carrega MFEs dinamicamente",
    "MFE Orders é o REMOTE: expõe OrdersApp como módulo independente via remoteEntry.js",
    "React e react-router-dom são singleton compartilhados — evita conflito de versões",
    "Cada MFE tem seu próprio Dockerfile e pode ser deployado sem rebuild do Shell",
], 0.5, 4.4, 12.3, 2.7, font_size=14, icon="▸")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — DEMONSTRAÇÃO AO VIVO (roteiro)
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
add_rect(sl, 0, 0, 13.33, 7.5, RGBColor(0xF4, 0xF7, 0xF6))
slide_header(sl, "Demonstração ao Vivo", "Roteiro — seguir esta ordem")

steps_demo = [
    ("🚀  Subir a stack",
     "docker-compose up --build\nAguardar ~2 min; todos os containers ficam healthy"),
    ("🌱  Seed de dados",
     "python scripts/seed.py\nCria admin@demo.com / admin123 + 3 pedidos de exemplo"),
    ("🔐  Login",
     "http://localhost:3000 → entrar com admin@demo.com\nObservar JWT salvo no localStorage (DevTools)"),
    ("📋  Listagem + Filtro",
     "Ver pedidos na tabela\nFiltrar por status (pending, processing, shipped...)"),
    ("➕  Criar Pedido",
     "Formulário dinâmico com itens\nAo salvar: IA sugere prioridade e gera resumo"),
    ("🔄  Atualizar Status",
     "Abrir detalhe → mudar status\nRedis cache é invalidado; lista mostra novo status"),
    ("💥  Resiliência",
     "docker-compose stop users-service\nOrders ainda lista; só PATCH /status retorna 401"),
]

cols = [(0.4, 1.7), (4.7, 1.7), (9.0, 1.7),
        (0.4, 4.35), (4.7, 4.35), (9.0, 4.35), (4.7, 4.35)]

positions = [
    (0.4, 1.65, 3.9), (4.65, 1.65, 3.9), (9.0, 1.65, 3.9),
    (0.4, 4.3, 3.9),  (4.65, 4.3, 3.9),  (9.0, 4.3, 3.9),
]

for i, ((titulo, desc), (left, top, w)) in enumerate(zip(steps_demo[:6], positions)):
    add_rect(sl, left, top, w, 2.45, BRANCO, VERDE_CLARO)
    add_text_box(sl, titulo, left + 0.15, top + 0.1, w - 0.3, 0.5,
                 font_size=13, bold=True, color=VERDE_GOV)
    add_text_box(sl, desc, left + 0.15, top + 0.65, w - 0.3, 1.65,
                 font_size=11, color=CINZA_ESCURO)

# step 7 (resiliência) em destaque no rodapé
add_rect(sl, 0.4, 6.85, 12.4, 0.5, RGBColor(0xFF, 0xF3, 0xCD), RGBColor(0xCC, 0x7A, 0x00))
add_text_box(sl, "💥  Resiliência: docker-compose stop users-service → Orders ainda funciona para leitura; "
             "apenas PATCH /status retorna 401 (sem JWT válido).",
             0.6, 6.88, 12.1, 0.45, font_size=11, color=RGBColor(0x7A, 0x4A, 0x00))


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — TOUR PELO CÓDIGO
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
add_rect(sl, 0, 0, 13.33, 7.5, RGBColor(0xF4, 0xF7, 0xF6))
slide_header(sl, "Tour pelo Código", "Separação de responsabilidades")

add_text_box(sl, "services/users/app/", 0.5, 1.65, 5.8, 0.4,
             font_size=13, bold=True, color=VERDE_GOV)
add_code_box(sl, """\
main.py       ← FastAPI app + lifespan (cria tabelas)
auth.py       ← bcrypt hash + JWT issuer (HS256)
models.py     ← User (UUID, email, role: admin|operator)
routes/
  auth.py     ← POST /auth/login
  users.py    ← CRUD /users""", 0.5, 2.1, 5.9, 2.1, font_size=11)

add_text_box(sl, "services/orders/app/", 6.9, 1.65, 5.9, 0.4,
             font_size=13, bold=True, color=VERDE_GOV)
add_code_box(sl, """\
main.py         ← FastAPI app
auth.py         ← JWT validator (shared secret)
ai_service.py   ← Claude Haiku: prioridade + resumo
redis_client.py ← Cache wrapper (graceful degrade)
routes/
  orders.py     ← CRUD + PATCH /status""", 6.9, 2.1, 5.9, 2.1, font_size=11)

add_text_box(sl, "frontend/", 0.5, 4.35, 12.3, 0.4,
             font_size=13, bold=True, color=VERDE_GOV)
add_code_box(sl, """\
shell/src/
  App.jsx              ← rotas protegidas + lazy import do MFE
  context/AuthContext  ← estado global do token JWT
  components/Login.jsx ← formulário de login

mfe-orders/src/
  OrdersApp.jsx        ← roteador interno do MFE
  api.js               ← fetch wrapper com Bearer header automático
  components/          ← OrdersList, OrderCreate, OrderDetail""",
             0.5, 4.8, 12.3, 2.3, font_size=11)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — INTEGRAÇÃO COM IA
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
add_rect(sl, 0, 0, 13.33, 7.5, RGBColor(0xF4, 0xF7, 0xF6))
slide_header(sl, "Integração com IA — Bônus", "Claude Haiku sugere prioridade e gera resumo do pedido")

add_code_box(sl, """\
# services/orders/app/ai_service.py
def suggest_priority_and_summary(customer_name: str, items: list, notes=None) -> dict:
    if not settings.ANTHROPIC_API_KEY:
        return {"priority": "medium", "summary": None}

    client = anthropic.Anthropic(api_key=settings.ANTHROPIC_API_KEY)
    prompt = "Analise o pedido e retorne APENAS JSON"

    message = client.messages.create(
        model="claude-haiku-4-5-20251001",
        max_tokens=200,
        messages=[{"role": "user", "content": prompt}]
    )
    result = json.loads(message.content[0].text.strip())
    return {"priority": result.get("priority", "medium"),
            "summary": result.get("summary")}""",
             0.5, 1.65, 12.3, 3.55, font_size=11)

add_bullet_box(sl, [
    "Chamado automaticamente ao POST /orders/ — resultado salvo em priority e ai_summary no banco",
    "Modelo Haiku escolhido por latência (~1s) e custo — adequado para classificação simples",
    "Fallback: se ANTHROPIC_API_KEY ausente, usa priority='medium' e ai_summary=None sem erro",
    "Em qualquer exceção na chamada da IA, o pedido é criado com fallback padrão",
], 0.5, 5.4, 12.3, 1.8, font_size=13, icon="▸")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — O QUE FOI ENTREGUE vs REQUISITOS
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
add_rect(sl, 0, 0, 13.33, 7.5, RGBColor(0xF4, 0xF7, 0xF6))
slide_header(sl, "O Que Foi Entregue", "Cobertura dos requisitos do edital")

rows_obrig = [
    ("2 microsserviços Python independentes",   True),
    ("PostgreSQL próprio por serviço",           True),
    ("Endpoints mínimos (listar/criar/status/ID/usuários)", True),
    ("MFE Shell + MFE Pedidos (listagem, filtro, criação)", True),
    ("Docker Compose (stack completa)",          True),
    ("README (instruções + diagrama + decisões)", True),
]
rows_desej = [
    ("Banco não-relacional (Redis)",             True),
    ("Testes automatizados (pytest)",            True),
    ("Pipeline de CI (GitHub Actions)",          True),
    ("Autenticação JWT compartilhada",           True),
]
rows_bonus = [
    ("Integração com IA (Claude Haiku)",         True),
    ("Documentação API (Swagger/OpenAPI)",       True),
    ("Observabilidade (logs estruturados)",      "parcial"),
    ("Comunicação assíncrona (filas)",           False),
]

def table_col(slide, titulo, rows, left, top, width):
    add_text_box(slide, titulo, left, top, width, 0.4,
                 font_size=13, bold=True, color=VERDE_GOV)
    for i, (label, status) in enumerate(rows):
        add_rect(slide, left, top + 0.45 + i * 0.46, width, 0.4,
                 BRANCO if i % 2 == 0 else RGBColor(0xF0, 0xF4, 0xF2))
        icon = "✅" if status is True else ("⚠️" if status == "parcial" else "❌")
        add_text_box(slide, f"{icon}  {label}",
                     left + 0.1, top + 0.48 + i * 0.46, width - 0.2, 0.38,
                     font_size=11, color=CINZA_ESCURO)

table_col(sl, "OBRIGATÓRIOS",  rows_obrig, 0.4,  1.65, 5.2)
table_col(sl, "DESEJÁVEIS",    rows_desej, 5.8,  1.65, 3.5)
table_col(sl, "BÔNUS",         rows_bonus, 9.5,  1.65, 3.5)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — DECISÕES NÃO TOMADAS
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
add_rect(sl, 0, 0, 13.33, 7.5, RGBColor(0xF4, 0xF7, 0xF6))
slide_header(sl, "Decisões Não Tomadas", "O que ficou de fora e por quê — isso é tão revelador quanto o código")

decisoes_nao = [
    ("API Gateway\n(Nginx/Kong)",
     "Atualmente o frontend chama cada serviço diretamente.\nCom mais tempo: roteamento centralizado, rate limiting, SSL."),
    ("MongoDB",
     "Optei por Redis — mais simples operacionalmente.\nMongoDB faria sentido se os pedidos tivessem schema variável."),
    ("JWKS / Validação\nRemota de Token",
     "Permite revogação imediata de tokens, mas adiciona\nlatência em cada request. Shared secret é suficiente no PMV."),
    ("Comunicação\nAssíncrona (filas)",
     "Publicar evento ao criar pedido (Redis Pub/Sub / RabbitMQ)\npermite notificações e auditoria sem acoplamento — Prioridade Média."),
    ("Refresh Token",
     "JWT atual expira em 24h sem mecanismo de refresh.\nUsuário precisa re-logar — aceitável para ferramenta interna."),
    ("RBAC real",
     "Qualquer autenticado pode atualizar status.\nCom mais tempo: regras de papéis (admin vs operator)."),
]

for i, (titulo, desc) in enumerate(decisoes_nao):
    col = 0.4 + (i % 3) * 4.3
    row = 1.7 + (i // 3) * 2.6
    add_rect(sl, col, row, 4.0, 2.3, BRANCO, AMARELO)
    add_text_box(sl, titulo, col + 0.15, row + 0.1, 3.7, 0.7,
                 font_size=14, bold=True, color=RGBColor(0x7A, 0x4A, 0x00))
    add_text_box(sl, desc, col + 0.15, row + 0.85, 3.7, 1.3,
                 font_size=11, color=CINZA_ESCURO)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — O QUE FARIA COM MAIS TEMPO
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
add_rect(sl, 0, 0, 13.33, 7.5, RGBColor(0xF4, 0xF7, 0xF6))
slide_header(sl, "O Que Faria com Mais Tempo", "Priorizados por impacto")

alta = [
    "API Gateway (Nginx/Kong): roteamento centralizado, rate limiting",
    "Refresh token: re-login sem expirar a sessão em 24h",
    "RBAC: regras de papéis (admin faz tudo; operator só leitura + criação)",
    "Testes de integração com PostgreSQL real (pytest-docker)",
]
media = [
    "Comunicação assíncrona: Redis Pub/Sub ao criar pedido",
    "MFE Catálogo: gerenciar produtos (hoje são strings livres)",
    "Paginação no frontend: crítico com milhares de pedidos",
    "Observabilidade: Prometheus + Grafana ou OpenTelemetry",
]

add_rect(sl, 0.4, 1.65, 5.9, 5.4, BRANCO, RGBColor(0xCC, 0x30, 0x2A))
add_text_box(sl, "🔴  Prioridade Alta", 0.6, 1.72, 5.5, 0.5,
             font_size=15, bold=True, color=RGBColor(0xCC, 0x30, 0x2A))
add_bullet_box(sl, alta, 0.6, 2.3, 5.5, 4.5, font_size=13, icon="▶")

add_rect(sl, 6.9, 1.65, 5.9, 5.4, BRANCO, AMARELO)
add_text_box(sl, "🟡  Prioridade Média", 7.1, 1.72, 5.5, 0.5,
             font_size=15, bold=True, color=RGBColor(0x7A, 0x4A, 0x00))
add_bullet_box(sl, media, 7.1, 2.3, 5.5, 4.5, font_size=13, icon="▶")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 14 — PERGUNTAS
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
add_rect(sl, 0, 0, 13.33, 7.5, VERDE_GOV)
add_rect(sl, 0, 5.8, 13.33, 1.7, RGBColor(0x0D, 0x3D, 0x2F))

add_text_box(sl, "Obrigado!", 1.0, 1.5, 11.33, 1.4,
             font_size=52, bold=True, color=BRANCO, align=PP_ALIGN.CENTER)
add_text_box(sl, "Perguntas do Avaliador", 1.0, 3.1, 11.33, 0.8,
             font_size=26, color=AMARELO, align=PP_ALIGN.CENTER)
add_text_box(sl, "Fique à vontade para aprofundar qualquer ponto da apresentação.",
             1.5, 4.0, 10.33, 0.6, font_size=16,
             color=RGBColor(0xB5, 0xE8, 0xD5), align=PP_ALIGN.CENTER)

add_text_box(sl, "Repositório · README · /docs (Swagger) · docker-compose up",
             1.0, 5.95, 11.33, 0.5,
             font_size=13, color=RGBColor(0x8A, 0xC9, 0xB5), align=PP_ALIGN.CENTER)


# ── Salvar ───────────────────────────────────────────────────────────────────
output = "Apresentacao_PMV_Gestao_Pedidos.pptx"
prs.save(output)
print(f"[OK] Arquivo salvo: {output}")
print(f"     Slides: {len(prs.slides)}")
